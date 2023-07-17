import collections
import collections.abc
c = collections
c.abc = collections.abc

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx import Presentation

# PPT 객체 생성
prs = Presentation()

def create_ppt():
    # 슬라이드 추가
    slide_layout = prs.slide_layouts[1]  # 레이아웃 선택 (1번은 제목 + 내용 레이아웃)
    slide = prs.slides.add_slide(slide_layout)

    # 텍스트 추가 및 스타일 변경
    text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2))
    text_frame = text_box.text_frame
    text_frame.text = "테스트"
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # 가운데 정렬
    text_frame.paragraphs[0].font.size = Pt(24)  # 글자 크기 변경
    text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)  # 글자 색 변경 (RGBColor 사용)

    # 이미지 추가 및 크기 변경
    slide_layout = prs.slide_layouts[6]  # 레이아웃 선택 (6번은 이미지 레이아웃)
    slide = prs.slides.add_slide(slide_layout)

    left = top = Inches(1)
    slide.shapes.add_picture('image.jpg', left, top, width=prs.slide_width, height=prs.slide_height)

    # 표 추가
    slide_layout = prs.slide_layouts[5]  # 레이아웃 선택 (5번은 두 콘텐츠가 있는 표 레이아웃)
    slide = prs.slides.add_slide(slide_layout)

    shapes = slide.shapes
    left = top = Inches(1)
    width = height = Inches(5)

    table = shapes.add_table(5, 3, left, top, width, height).table

    # 표 스타일 설정
    table.cell(0, 0).fill.solid()  # 타이틀 셀 배경색 설정
    table.cell(0, 0).fill.fore_color.rgb = RGBColor(255, 255, 240)  # ivory색 (RGBColor 사용)

    # 표 셀에 데이터 추가
    table.cell(0, 0).text = "Title 1"
    table.cell(0, 1).text = "Title 2"
    table.cell(0, 2).text = "Title 3"

    for row in range(1, 5):
        for col in range(3):
            table.cell(row, col).text = "Cell {},{}".format(row, col)

    # PPT 파일 저장
    prs.save("output.pptx")
